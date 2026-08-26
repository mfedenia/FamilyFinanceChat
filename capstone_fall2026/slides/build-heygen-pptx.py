#!/usr/bin/env python3
"""
Build the HeyGen-ready PowerPoint from the build-step Beamer deck.

Pipeline
--------
    familyfinancechat-fall2026-build.tex   35 cumulative overlay steps
        -> latexmk                          familyfinancechat-fall2026-build.pdf (35 pages)
        -> pdftoppm                         one PNG per page, 1920x1080
        -> python-pptx                      familyfinancechat-fall2026-heygen.pptx

Each PDF page becomes one 16:9 PowerPoint slide holding that page as a
full-bleed image, with its narration line in the speaker notes.  HeyGen reads
the speaker notes as the script for that slide, so the spoken words always match
what has just appeared on screen.

NARRATION is the single source of truth for the video script.  It must have
exactly one entry per PDF page, in order -- the script refuses to build if the
counts disagree, which is what stops the deck and the speech drifting apart.

This is a SHORTER script than ../script/presentation-script.md (which is the
five-minute version for a live human presenter).  Target here is under four
minutes.

Requires: latexmk + pdflatex, pdftoppm (poppler), python-pptx.
    python3 -m venv venv && venv/bin/pip install python-pptx
    venv/bin/python build-heygen-pptx.py
"""

import shutil
import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
TEX = HERE / "familyfinancechat-fall2026-build.tex"
PDF = HERE / "familyfinancechat-fall2026-build.pdf"
PNG_DIR = HERE / "build-png"
PPTX = HERE / "familyfinancechat-fall2026-heygen.pptx"

# Words per minute a HeyGen avatar speaks at, for the runtime estimate below.
# Deliberately the slow end of the plausible range, so the estimate errs long.
WPM = 150

# ---------------------------------------------------------------------------
# One entry per build step, in page order.  (deck_slide, step, narration)
# deck_slide/step are documentation only -- the order of this list is what maps
# narration onto pages.
#
# VOICE: this is somebody talking, not a slide read aloud.  The deck already
# carries the facts; the narration carries the reason to care about them.  So:
# whole sentences with connective tissue ("so", "because", "here's why"), a
# concrete picture before an abstraction, and the stakes said out loud.  Where
# the deck has a fragment, the speech has a thought.
#
# Contractions are deliberate.  "We're" and "doesn't" are what a person says;
# "we are" and "does not" are what a document says, and an avatar reading the
# document version sounds like a kiosk.
#
# WRITTEN FOR THE EAR, NOT THE EYE.  Every one of these is spoken aloud by a
# text-to-speech engine that has no idea what any of it means, so the script is
# spelled the way it should sound, not the way it is written elsewhere:
#
#   numbers as words        "one point two seconds", not "1.2s" -- a speech
#                           engine reads "1.2" as "one point two" if you are
#                           lucky and as a date or version number if you are not
#   acronyms letter by letter   "A-I", "U-I" -- hyphens force the letters apart.
#                           Written solid, "AI" and "UI" get guessed at as words
#                           ("ay", "ooey").  Same trick for any acronym added
#                           later: S-Q-L, A-P-I, U-R-L.
#   product names as spoken     "Family Finance Chat", "Open Web U-I" -- run-together
#                           capitals (FamilyFinanceChat, WebUI) are read as one
#                           invented word.  The SLIDES still show the real
#                           spelling; only the speech is respelled.
#   dates spoken in full    "the twentieth of October", "the middle of November"
#                           -- never "20 Oct" or "mid-Nov".
#   no em dashes, no semicolons     a full stop or a comma is an unambiguous
#                           pause.  An em dash is read as anything from a beat
#                           to nothing at all, depending on the engine.
#
# check_tts() below enforces the mechanical half of this on every build.
# ---------------------------------------------------------------------------
NARRATION = [
    (1, 1, "This is Family Finance Chat. Two jobs this semester: make it solid, and "
           "give it a face."),

    (2, 1, "Picture a student practising a hard conversation about money. Across from "
           "them, an A-I family built from the course documents. It never invents facts."),
    (2, 2, "Every question is scored, and the professor sees it all in one place."),
    (2, 3, "And that's the real gift. Practice at midnight, a hundred times, with "
           "nobody's schedule to book."),

    (3, 1, "It works. Students use it today. But three things hold it back. We're three "
           "versions behind the software underneath, and the gap grows monthly."),
    (3, 2, "The professor's grading tool only runs on a laptop, and takes a terminal and "
           "a developer."),
    (3, 3, "And every deployment leans on steps someone has to remember. Sometimes they don't."),
    (3, 4, "None of it's glamorous. But it's the distance between a system that works and "
           "one you can hand over."),

    (4, 1, "One rule shapes everything: we never touch the core of Open Web U-I, the "
           "platform underneath."),
    (4, 2, "The last team learned the hard way. Six and a half thousand lines of copied "
           "code, frozen on one old version. They deleted it all. Ours is one line."),
    (4, 3, "Before we build anything we ask: if they shipped a release tomorrow, would we "
           "still work?"),

    (5, 1, "The first job is what the course depends on. Upgrade to current, one careful "
           "step at a time. These rewrite the database. There's no undo."),
    (5, 2, "Give the professor a web address. No laptop, no terminal, no developer."),
    (5, 3, "Automate the testing, so a bad change is caught by us, not a student."),
    (5, 4, "And clear out the last of that old fork, so nobody drags it back."),

    (6, 1, "Now the part we're excited about. Today, the student types."),
    (6, 2, "We want them to speak. To a face that listens and answers out loud, in character."),
    (6, 3, "Because soon they'll sit with a real family, asking the hardest question out "
           "loud, once, while everyone watches. You can't rehearse that with a backspace key."),

    (7, 1, "So how do we add a face without breaking our rule?"),
    (7, 2, "We don't build it inside the platform. We build it alongside. Its own container."),
    (7, 3, "The student talks to it, it asks Open Web U-I for the answer, and the platform "
           "never notices. Same brain, same grading."),
    (7, 4, "And if it doesn't pan out, we delete one container. The rule didn't limit the "
           "design. It gave us a better one."),

    (8, 1, "Three numbers tell us whether this is real. Speed. It answers in about one "
           "point two seconds, or it stops feeling like a conversation."),
    (8, 2, "Cost. Around four hundred dollars a semester for a class of forty, capped so "
           "it can't surprise us."),
    (8, 3, "And permission. A student's voice and face are protected records. Consent is "
           "settled in week two."),
    (8, 4, "If any one fails, we stop. And stopping is a perfectly good answer."),

    (9, 1, "By late September, students can talk to it. Voice only, no face yet."),
    (9, 2, "By the twentieth of October, we pick an avatar provider, or decide not to."),
    (9, 3, "A week later, the platform is current."),
    (9, 4, "And by the middle of November, five students hold full spoken sessions, graded "
           "like everyone else."),
    (9, 5, "The voice version ships in September no matter what. That's what makes the "
           "ambitious half safe."),

    (10, 1, "So picture December. A professor grading from a web address. No laptop, no help."),
    (10, 2, "Deployments that run themselves, so a bad change never reaches a student."),
    (10, 3, "And real evidence about whether talking to a face makes better advisors. It "
            "might not. We publish either way."),
    (10, 4, "Because we're not here just to finish a system. We're here to leave one where "
            "the interesting work is the teaching, not the plumbing."),
]


def check_tts():
    """Catch the spellings a speech engine reliably gets wrong.

    Mechanical only -- it cannot tell you the writing is stiff, just that a digit
    or a run-together acronym slipped in.  Fatal, not a warning: a mispronounced
    word is not something you notice until the video is rendered and paid for.
    """
    import re
    rules = [
        (r"\d",                 "digits — write the number in words"),
        (r"[A-Z][a-z]+[A-Z]",   "run-together capitals — respell as separate spoken words"),
        (r"\b[A-Z]{2,}\b",      "solid acronym — hyphenate it, like A-I or S-Q-L"),
        (r"—",                  "em dash — use a full stop or a comma"),
        (r";",                  "semicolon — use a full stop"),
        (r"\b\w+-(?![A-Z]\b)\w",  "hyphenated word — say it in full, like 'the middle of November'"),
    ]
    problems = [f"     slide {d}.{s}: {why}\n       {n}"
                for d, s, n in NARRATION
                for pat, why in rules if re.search(pat, n)]
    if problems:
        sys.exit("NARRATION is not safe for text-to-speech:\n" + "\n".join(problems))


def run(cmd, **kw):
    proc = subprocess.run(cmd, capture_output=True, text=True, **kw)
    if proc.returncode != 0:
        sys.exit(f"FAILED: {' '.join(cmd)}\n{proc.stdout[-2000:]}\n{proc.stderr[-2000:]}")
    return proc


def build_pdf():
    print("1/4  compiling the build-step deck")
    run(["latexmk", "-pdf", "-interaction=nonstopmode", "-halt-on-error", TEX.name], cwd=HERE)


def render_pngs():
    print("2/4  rendering pages to 1920x1080 PNG")
    if PNG_DIR.exists():
        shutil.rmtree(PNG_DIR)
    PNG_DIR.mkdir()
    # -scale-to-x/-scale-to-y forces exactly 1920x1080; the deck is already 16:9,
    # so nothing is distorted, and HeyGen gets a clean full-HD frame.
    run(["pdftoppm", "-png", "-r", "300", "-scale-to-x", "1920", "-scale-to-y", "1080",
         str(PDF), str(PNG_DIR / "slide")])
    return sorted(PNG_DIR.glob("slide-*.png"))


def build_pptx(pngs):
    print("3/4  assembling the PowerPoint")
    from pptx import Presentation
    from pptx.util import Emu, Inches

    if len(pngs) != len(NARRATION):
        sys.exit(f"MISMATCH: {len(pngs)} PDF pages but {len(NARRATION)} narration entries.\n"
                 f"Every build step needs exactly one line of narration — fix NARRATION "
                 f"in this file, or the overlay specs in {TEX.name}.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for png, (deck_slide, step, narration) in zip(pngs, NARRATION):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), Emu(0), Emu(0),
                                 width=prs.slide_width, height=prs.slide_height)
        # HeyGen reads the speaker notes as the script for this slide.
        slide.notes_slide.notes_text_frame.text = narration

    prs.save(PPTX)


def report():
    print("4/4  done")
    words = sum(len(n.split()) for _, _, n in NARRATION)
    seconds = words / WPM * 60
    print(f"\n     {PPTX.name}")
    print(f"     {len(NARRATION)} slides · {words} words · about "
          f"{int(seconds // 60)}:{int(seconds % 60):02d} at {WPM} words per minute")
    if seconds > 240:
        print("     WARNING: over the four-minute target — trim NARRATION.")


if __name__ == "__main__":
    check_tts()
    build_pdf()
    pngs = render_pngs()
    build_pptx(pngs)
    report()
